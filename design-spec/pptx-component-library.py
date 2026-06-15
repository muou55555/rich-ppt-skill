"""
rich-pptx Component Library
源自 GenSpark HTML Deck 设计语言分析

用法：在生成PPT的Python脚本中 import 或复制相关函数
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ─── 命名空间 ──────────────────────────────────────────────────────────────

_NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

def qn(tag):
    prefix, local = tag.split(':')
    return f'{{{_NS[prefix]}}}{local}'

# ─── 单位换算 ──────────────────────────────────────────────────────────────

def E(inch):
    """英寸 → EMU"""
    return int(inch * 914400)

def PX(px_val):
    """1280px 画布像素 → 英寸（13.333"/1280px）"""
    return px_val * 13.333 / 1280

def PY(px_val):
    """720px 画布像素 → 英寸（7.5"/720px）"""
    return px_val * 7.5 / 720

def px_to_pt(px_val):
    """屏幕像素(96dpi) → 点数"""
    return px_val * 72 / 96

# ─── 配色系统 ──────────────────────────────────────────────────────────────

C = {
    # 背景层
    'bg':       RGBColor(0x0F, 0x17, 0x2A),  # slate-900 主背景
    'card':     RGBColor(0x1E, 0x29, 0x3B),  # slate-800 卡片背景
    'panel':    RGBColor(0x1E, 0x29, 0x3B),  # slate-800 左侧面板
    'border':   RGBColor(0x33, 0x41, 0x55),  # slate-700 边框/分隔
    # 强调色
    'accent':   RGBColor(0x38, 0xBD, 0xF8),  # sky-400  装饰线/图标
    'accent2':  RGBColor(0x0E, 0xA5, 0xE9),  # sky-500  渐变/发光
    'blue':     RGBColor(0x3B, 0x82, 0xF6),  # blue-500 链接/按钮
    # 文字层
    'text':     RGBColor(0xF1, 0xF5, 0xF9),  # slate-100 主标题
    'text2':    RGBColor(0xCB, 0xD5, 0xE1),  # slate-300 副标题
    'muted':    RGBColor(0x94, 0xA3, 0xB8),  # slate-400 正文
    'faint':    RGBColor(0x64, 0x74, 0x8B),  # slate-500 弱化
    # 语义色
    'danger':   RGBColor(0xEF, 0x44, 0x44),  # red-500
    'warning':  RGBColor(0xF5, 0x9E, 0x0B),  # amber-500
    'success':  RGBColor(0x22, 0xC5, 0x5E),  # green-500
    'purple':   RGBColor(0xA8, 0x55, 0xF7),  # violet-500
}

# ─── 基础图形操作 ──────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, border=None, bpt=0.75):
    """
    添加矩形
    x/y/w/h: 英寸
    fill: RGBColor 或 None（透明）
    border: RGBColor 或 None（无边框）
    """
    s = slide.shapes.add_shape(1, E(x), E(y), E(w), E(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(bpt)
    else:
        s.line.fill.background()
    return s


def add_oval(slide, x, y, w, h, fill=None, alpha_pct=100):
    """
    添加椭圆（用于装饰圆/发光圆）
    alpha_pct: 0-100 透明度百分比（100=不透明）
    """
    s = slide.shapes.add_shape(9, E(x), E(y), E(w), E(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        if alpha_pct < 100:
            # 通过 XML 设置透明度
            sp = s._element
            solidFill = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
            if solidFill is not None:
                srgbClr = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                if srgbClr is not None:
                    alpha = etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                    # alpha val: 100000=不透明, 0=全透明
                    alpha.set('val', str(int(alpha_pct * 1000)))
    else:
        s.fill.background()
    s.line.fill.background()
    return s


def round_(shape, adj=16667):
    """
    将矩形转为圆角矩形
    adj: 圆角程度 0(直角) ~ 50000(pill)
    常用值: 5000=微小圆角, 8000=卡片, 16667≈10%, 50000=pill
    """
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    pg = spPr.find(qn('a:prstGeom'))
    pg.set('prst', 'roundRect')
    al = pg.find(qn('a:avLst'))
    if al is None:
        al = etree.SubElement(pg, qn('a:avLst'))
    for g in list(al):
        al.remove(g)
    gd = etree.SubElement(al, qn('a:gd'))
    gd.set('name', 'adj')
    gd.set('fmla', f'val {adj}')

# ─── 文本操作 ──────────────────────────────────────────────────────────────

def add_tb(slide, x, y, w, h, lines, size, bold=False,
           color=None, align=PP_ALIGN.LEFT, ls=1.15, font='Noto Sans SC'):
    """
    添加多行文本框
    lines: 字符串列表，每个元素一段
    size: 字体 pt 大小
    ls: 行距倍数（相对于字体大小）
    """
    txBox = slide.shapes.add_textbox(E(x), E(y), E(w), E(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = Pt(size * ls)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or C['text']
        run.font.name = font
    return txBox


def add_tb_multicolor(slide, x, y, w, h, pairs, size, ls=1.2, font='Noto Sans SC'):
    """
    多色文本框：每行可指定独立颜色
    pairs: [(text, RGBColor), ...]
    """
    txBox = slide.shapes.add_textbox(E(x), E(y), E(w), E(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, color) in enumerate(pairs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = Pt(size * ls)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font
    return txBox

# ─── 复合组件 ──────────────────────────────────────────────────────────────

def set_bg(slide, color=None):
    """设置幻灯片背景色，默认 GenSpark 深色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color or C['bg']


def add_deco_line(slide, x, y, w=0.625):
    """
    标题上方的蓝色装饰横线
    默认宽 0.625"（60px）
    """
    r = add_rect(slide, x, y, w, 0.042, fill=C['accent'])
    return r


def add_decor_circles(slide, size=3.125):
    """
    右上角和左下角半透明装饰圆（通用装饰）
    size: 圆的直径（英寸），默认 3.125" ≈ 300px
    """
    # 右上角
    tl = add_oval(slide, 13.333 - size * 0.6, -size * 0.6, size, size,
                  fill=C['accent'], alpha_pct=5)
    # 左下角
    br = add_oval(slide, -size * 0.6, 7.5 - size * 0.4, size, size,
                  fill=C['accent'], alpha_pct=5)
    return tl, br


def add_bg_glow(slide):
    """发光装饰圆（巨大椭圆 + 极低透明度）"""
    glow_size = 6.25  # 600px → 6.25"
    # 右上发光
    g1 = add_oval(slide, 13.333 - glow_size * 0.7, -glow_size * 0.7,
                  glow_size, glow_size, fill=C['accent2'], alpha_pct=5)
    # 左下发光
    g2 = add_oval(slide, -glow_size * 0.7, 7.5 - glow_size * 0.3,
                  glow_size, glow_size, fill=C['accent2'], alpha_pct=5)
    return g1, g2


def add_feature_card(slide, x, y, w, h,
                     icon_char, icon_color, icon_bg,
                     title, body,
                     title_size=11, body_size=9):
    """
    带图标的功能卡片
    icon_char: emoji 字符或 Unicode 图标
    """
    # 卡片背景
    card = add_rect(slide, x, y, w, h, fill=C['card'], border=C['border'], bpt=0.5)
    round_(card, adj=8000)

    icon_w = 0.35
    pad = 0.12
    gap = 0.1

    # 图标盒
    ib = add_rect(slide, x + pad, y + (h - icon_w) / 2,
                  icon_w, icon_w, fill=icon_bg)
    round_(ib, adj=8000)
    add_tb(slide, x + pad, y + (h - icon_w) / 2,
           icon_w, icon_w, [icon_char], title_size,
           color=icon_color, align=PP_ALIGN.CENTER)

    # 文字区域
    tx = x + pad + icon_w + gap
    tw = w - pad - icon_w - gap - pad * 0.5
    th_title = h * 0.45
    th_body = h * 0.45

    add_tb(slide, tx, y + h * 0.1, tw, th_title,
           [title], title_size, bold=True, color=C['text'])
    add_tb(slide, tx, y + h * 0.48, tw, th_body,
           [body], body_size, color=C['muted'])

    return card


def add_badge(slide, x, y, w, h, text, text_color=None, bg_color=None, size=8):
    """
    小型胶囊标签（pill badge）
    """
    bg = bg_color or C['card']
    tc = text_color or C['accent']
    b = add_rect(slide, x, y, w, h, fill=bg, border=tc, bpt=0.3)
    round_(b, adj=50000)
    add_tb(slide, x, y, w, h, [text], size, color=tc,
           align=PP_ALIGN.CENTER)
    return b


def add_header(slide, title, subtitle, badge_text=None,
               x=0.625, y=0.312, w=12.083):
    """
    通用页面 header（装饰线 + 标题 + 副标题）
    对应 GenSpark .header 区块
    """
    # 装饰线
    add_deco_line(slide, x, y)

    # 主标题
    add_tb(slide, x, y + 0.083, w, 0.42,
           [title], 27, bold=True, color=C['text'])

    # 副标题
    if subtitle:
        add_tb(slide, x, y + 0.5, w, 0.25,
               [subtitle], 12, color=C['muted'])

    # 可选章节 badge（右上角）
    if badge_text:
        bw = len(badge_text) * 0.1 + 0.25
        add_badge(slide, 13.333 - bw - 0.3, y + 0.1, bw, 0.22,
                  badge_text, text_color=C['accent'], bg_color=C['card'])


# ─── 布局模板 ──────────────────────────────────────────────────────────────

def layout_card_grid(slide, cards_data, start_x=0.625, start_y=1.2,
                     cols=2, col_gap=0.167, row_gap=0.167):
    """
    生成卡片网格
    cards_data: [(icon, icon_color, icon_bg, title, body), ...]
    """
    available_w = 13.333 - start_x * 2
    available_h = 7.5 - start_y - 0.2
    n = len(cards_data)
    rows = (n + cols - 1) // cols

    card_w = (available_w - col_gap * (cols - 1)) / cols
    card_h = (available_h - row_gap * (rows - 1)) / rows

    for i, data in enumerate(cards_data):
        col = i % cols
        row = i // cols
        cx = start_x + col * (card_w + col_gap)
        cy = start_y + row * (card_h + row_gap)
        add_feature_card(slide, cx, cy, card_w, card_h, *data)


def layout_split(slide, left_w_ratio=0.30, padding=0.625):
    """
    返回左右分栏的坐标参数
    left_w_ratio: 左侧占比（默认 30%）
    Returns: (left_x, left_y, left_w, left_h, right_x, right_y, right_w, right_h)
    """
    total_w = 13.333 - padding * 2
    gap = 0.25
    lw = total_w * left_w_ratio
    rw = total_w * (1 - left_w_ratio) - gap
    lx = padding
    rx = padding + lw + gap
    y = 1.1
    h = 7.5 - y - 0.2
    return lx, y, lw, h, rx, y, rw, h


# ─── 完整页面模板 ──────────────────────────────────────────────────────────

def make_dark_slide(prs):
    """创建一张深色背景的空白幻灯片"""
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_decor_circles(slide)
    add_bg_glow(slide)
    return slide


def make_content_slide(prs, title, subtitle, badge=None):
    """
    创建带 header 的内容页
    返回 slide 对象，供后续添加内容
    """
    slide = make_dark_slide(prs)
    add_header(slide, title, subtitle, badge_text=badge)
    return slide


# ─── 示例：4卡片内容页 ─────────────────────────────────────────────────────

if __name__ == '__main__':
    prs = Presentation()
    prs.slide_width = E(13.333)
    prs.slide_height = E(7.5)

    slide = make_content_slide(prs, 'Claude Code 核心特性', '基于多智能体架构的终端原生 AI 编程助手', '核心架构')

    cards = [
        ('🤖', C['accent'], RGBColor(0x0A, 0x1F, 0x3A), '多智能体编排', '支持并行 subagent、任务分解与上下文共享'),
        ('🔀', C['accent'], RGBColor(0x0A, 0x1F, 0x3A), '工具生态集成', 'Bash/Read/Write/WebFetch 等 30+ 内置工具'),
        ('🛡️', C['accent'], RGBColor(0x0A, 0x1F, 0x3A), '权限安全模型', '最小权限原则，危险操作需明确授权'),
        ('⚡', C['accent'], RGBColor(0x0A, 0x1F, 0x3A), '流式响应', '实时输出 token，低延迟交互体验'),
    ]
    layout_card_grid(slide, cards, cols=2)

    out_path = '/Users/zhangxiaowu/Claude/Projects/rich-ppt-skill/design-spec/example_output.pptx'
    prs.save(out_path)
    print(f'✅ 示例 PPT 已生成：{out_path}')
