#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
components.py — Rich PPT 组件库
基于 python-pptx，提供圆角卡片、箭头流程、表格、标注框等高层组件。
"""
import os, sys, copy
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── 幻灯片尺寸 ──────────────────────────────────────────────────────────────
W = 12192000   # EMU 宽（宽屏 16:9）
H = 6858000    # EMU 高

# ── 色板 ─────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x00, 0x32, 0x9D)
MED    = RGBColor(0x30, 0x55, 0x98)
SKY    = RGBColor(0x44, 0x72, 0xC4)
STEEL  = RGBColor(0x2E, 0x75, 0xB6)
GREEN  = RGBColor(0x1E, 0x8C, 0x55)
GOLD   = RGBColor(0xED, 0xAD, 0x1A)
ORANGE = RGBColor(0xE0, 0x6B, 0x1A)
RED2   = RGBColor(0xC0, 0x39, 0x2B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CARD   = RGBColor(0xEB, 0xF0, 0xF8)
LINE   = RGBColor(0xCC, 0xD8, 0xEE)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
GRAY   = RGBColor(0x66, 0x72, 0x82)
ACCENTS = [NAVY, SKY, GREEN, GOLD, ORANGE, MED]


def _hex(c: RGBColor) -> str:
    """RGBColor → '#RRGGBB'"""
    return f"#{str(c)}"


# ── 基础绘图 ──────────────────────────────────────────────────────────────────

def new_presentation() -> Presentation:
    """新建 16:9 宽屏演示文稿"""
    prs = Presentation()
    prs.slide_width  = Emu(W)
    prs.slide_height = Emu(H)
    return prs


def blank_slide(prs: Presentation):
    """添加空白幻灯片（使用空白 layout）"""
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, bg=NAVY, r=0):
    """绘制矩形（r>0 时为圆角，r 单位为 EMU）"""
    shape_type = 5 if r > 0 else 1
    sh = slide.shapes.add_shape(shape_type, Emu(x), Emu(y), Emu(w), Emu(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = bg
    sh.line.fill.background()
    if r > 0:
        try:
            prstGeom = sh._element.spPr.prstGeom
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {r}')
        except Exception:
            pass
    return sh


def oval(slide, x, y, w, h, bg=NAVY):
    sh = slide.shapes.add_shape(9, Emu(x), Emu(y), Emu(w), Emu(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = bg
    sh.line.fill.background()
    return sh


def line_h(slide, x, y, w, color=LINE):
    sh = slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(w), Emu(int(H * 0.003)))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()


def box(slide, x, y, w, h, text, size=12, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color if color else DARK
    return txb


def multiline(slide, lines, x, y, w, h, size=11, color=None, gap=3, bold_first=False):
    txb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color if color else DARK
        run.font.bold = (bold_first and i == 0)
    return txb


def add_icon(slide, name: str, x: int, y: int, size_emu: int, color: str = None):
    """在幻灯片上插入 Phosphor fill 图标（PNG）"""
    try:
        from scripts.icons import get_icon
    except ImportError:
        sys.path.insert(0, os.path.dirname(__file__))
        from icons import get_icon
    hex_color = color if color else "#00329D"
    px = max(48, int(size_emu / 914400 * 96))
    icon_path = get_icon(name, hex_color, px)
    slide.shapes.add_picture(icon_path, Emu(x), Emu(y), Emu(size_emu), Emu(size_emu))


# ── 高级组件 ──────────────────────────────────────────────────────────────────

def slide_bg(slide, color=WHITE):
    """填充幻灯片背景"""
    rect(slide, 0, 0, W, H, color)


def title_block(slide, title: str, subtitle: str = ''):
    """
    标准标题块：左侧竖条 + 标题文字（NAVY 28pt）+ 水平线 + 副标题
    与 GordenPPTSkill 模板视觉语言对齐。
    """
    rect(slide, int(W * 0.025), int(H * 0.025), int(W * 0.006), int(H * 0.115), NAVY)
    box(slide, int(W * 0.045), int(H * 0.025), int(W * 0.88), int(H * 0.11),
        title, size=28, bold=False, color=NAVY)
    line_h(slide, int(W * 0.025), int(H * 0.145), int(W * 0.95), LINE)
    if subtitle:
        box(slide, int(W * 0.045), int(H * 0.155), int(W * 0.9), int(H * 0.06),
            subtitle, size=11, color=GRAY)


def num_badge(slide, x: int, y: int, n: int, bg=NAVY, size: int = 14):
    """圆形数字徽章"""
    d = int(H * 0.072)
    oval(slide, x, y, d, d, bg)
    box(slide, x, y, d, d, str(n), size=size, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)


def rcard(slide, x: int, y: int, w: int, h: int,
          title: str, body_lines, accent=NAVY, num=None, icon: str = None):
    """
    圆角内容卡片。
    body_lines: str 或 list[str]
    icon: 图标名（Phosphor 语义名）
    num: 编号数字（和 icon 二选一）
    """
    rect(slide, x, y, w, h, WHITE, r=20000)
    rect(slide, x, y, w, int(H * 0.007), accent)
    bdr = slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(w), Emu(h))
    bdr.fill.background()
    bdr.line.color.rgb = LINE
    bdr.line.width = Pt(0.5)

    xpad = int(W * 0.012)
    ty   = y + int(H * 0.018)
    icon_sz = int(H * 0.065)

    if icon:
        add_icon(slide, icon, x + xpad, ty, icon_sz, _hex(accent))
        tx = x + xpad + icon_sz + int(W * 0.008)
    elif num is not None:
        num_badge(slide, x + xpad, ty - int(H * 0.002), num, accent, 12)
        tx = x + xpad + int(H * 0.076)
    else:
        tx = x + xpad

    box(slide, tx, ty, w - (tx - x) - xpad, int(H * 0.06),
        title, size=13, bold=True, color=accent)

    lines = body_lines if isinstance(body_lines, list) else body_lines.split('\n')
    multiline(slide, lines,
              x + xpad, ty + int(H * 0.072),
              w - 2 * xpad, h - int(H * 0.11),
              size=10, color=DARK, gap=2)


def callout_box(slide, x: int, y: int, w: int, h: int,
                text: str, size: int = 10,
                border_color=NAVY, bg=None, icon: str = None):
    """标注框：左边框 + 浅色背景 + 可选图标"""
    if bg is None:
        bg = CARD
    rect(slide, x, y, w, h, bg, r=10000)
    rect(slide, x, y, int(W * 0.006), h, border_color)
    tx = x + int(W * 0.015)
    if icon:
        icon_sz = int(H * 0.055)
        add_icon(slide, icon, tx, y + int(H * 0.012), icon_sz, _hex(border_color))
        tx += icon_sz + int(W * 0.010)
    box(slide, tx, y + int(H * 0.01), w - (tx - x) - int(W * 0.015),
        h - int(H * 0.015), text, size=size, color=DARK)


def section_tag(slide, x: int, y: int, text: str, bg=NAVY, icon: str = None):
    """圆角胶囊标签，可选左侧图标"""
    tw = int(W * 0.16) if icon else int(W * 0.14)
    th = int(H * 0.048)
    rect(slide, x, y, tw, th, bg, r=15000)
    ix = x + int(W * 0.008)
    if icon:
        ic_sz = int(th * 0.7)
        add_icon(slide, icon, ix, y + (th - ic_sz) // 2, ic_sz, "#FFFFFF")
        ix += ic_sz + int(W * 0.005)
    box(slide, ix, y + int(H * 0.005), tw - (ix - x) - int(W * 0.005),
        th - int(H * 0.008), text, size=10, bold=True, color=WHITE)


def htable(slide, headers: list, rows: list,
           x: int, y: int, w: int, h: int, hbg=NAVY):
    """高亮表格：表头 hbg 底白字，奇偶行交替"""
    nc, nr = len(headers), len(rows) + 1
    tbl = slide.shapes.add_table(nr, nc, Emu(x), Emu(y), Emu(w), Emu(h)).table
    for ci, hdr in enumerate(headers):
        c = tbl.cell(0, ci)
        c.text = hdr
        c.fill.solid()
        c.fill.fore_color.rgb = hbg
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0] if p.runs else p.add_run()
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = WHITE
    for ri, row in enumerate(rows):
        bg = CARD if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            c = tbl.cell(ri + 1, ci)
            c.text = str(val)
            c.fill.solid()
            c.fill.fore_color.rgb = bg
            p = c.text_frame.paragraphs[0]
            r = p.runs[0] if p.runs else p.add_run()
            r.font.size = Pt(9)
            r.font.color.rgb = DARK
    return tbl


def chevron_flow(slide, steps: list, y: int, item_h: int, colors=None):
    """
    箭头流程图（横排 RIGHT_ARROW）
    steps: [(label, subtitle), ...]
    """
    n = len(steps)
    gap = int(W * 0.005)
    total_w = W - 2 * int(W * 0.03)
    item_w = (total_w - gap * (n - 1)) // n
    x0 = int(W * 0.03)
    for i, (lbl, sub) in enumerate(steps):
        x = x0 + i * (item_w + gap)
        clr = colors[i % len(colors)] if colors else ACCENTS[i % len(ACCENTS)]
        sh = slide.shapes.add_shape(13, Emu(x), Emu(y), Emu(item_w), Emu(item_h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = clr
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = lbl
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = WHITE
        if sub:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = sub
            r2.font.size = Pt(9)
            r2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)


def purge_orphans(prs: Presentation):
    """清理孤儿 slide 关系（操作 _sldIdLst 后必须调用）"""
    RT = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide'
    keep = {sl.get(qn('r:id')) for sl in prs.slides._sldIdLst}
    ppart = prs.part
    for rid in list(ppart.rels.keys()):
        if ppart.rels[rid].reltype == RT and rid not in keep:
            ppart.rels._rels.pop(rid)


def save_pptx(prs: Presentation, output_path: str):
    """保存演示文稿，自动创建目录"""
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    purge_orphans(prs)
    prs.save(output_path)
    size_kb = os.path.getsize(output_path) // 1024
    print(f"✓ Saved: {output_path}  ({prs.slides.__len__()} slides, {size_kb}KB)")
    return output_path


# ── 带媒体关系的模板页复制 ────────────────────────────────────────────────────
# 修复：原来只复制 spTree XML，丢失了背景图 / 嵌入图片的 relationship。
# 参考 references/hybrid-mode.md 获取完整说明。

def copy_template_slide(prs: Presentation, src_slide) -> object:
    """
    将 src_slide（来自另一个 Presentation）完整复制为 prs 中的新幻灯片。
    正确传递所有图片/媒体 relationship，避免图像丢失或幻灯片无法正常显示。

    用法：
        src_prs = Presentation("template.pptx")
        for idx in [0, 3, 5]:
            copy_template_slide(prs, src_prs.slides[idx])
    """
    from lxml import etree as _etree
    import copy as _copy

    blank = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank)
    new_part  = new_slide.part
    src_part  = src_slide.part

    # 1. 复制所有非 slideLayout relationship（图片、视频、音频等）
    rId_map = {}
    for rId, rel in list(src_part.rels.items()):
        if 'slideLayout' in rel.reltype:
            continue
        try:
            if rel.is_external:
                new_rId = new_part.relate_to(rel.target_ref, rel.reltype, is_external=True)
            else:
                new_rId = new_part.relate_to(rel.target_part, rel.reltype)
            rId_map[rId] = new_rId
        except Exception as e:
            print(f"  [warn] skip rel {rId}: {e}")

    # 2. 深复制整个 cSld（含背景 p:bg 和形状 p:spTree）
    src_cSld = _copy.deepcopy(src_slide._element.find(qn('p:cSld')))

    # 3. 将 XML 中旧 rId 替换为新 rId
    if rId_map:
        xml_str = _etree.tostring(src_cSld, encoding='unicode')
        for old_rId, new_rId in rId_map.items():
            if old_rId != new_rId:
                xml_str = xml_str.replace(f'"{old_rId}"', f'"{new_rId}"')
        src_cSld = _etree.fromstring(xml_str)

    # 4. 替换 new_slide 的 cSld
    dst = new_slide._element
    old_cSld = dst.find(qn('p:cSld'))
    if old_cSld is not None:
        dst.remove(old_cSld)
    dst.append(src_cSld)

    return new_slide
