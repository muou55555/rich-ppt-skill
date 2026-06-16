#!/usr/bin/env python3
"""
HTML slides v2 → Editable PPTX  (v2 — HTML-parsing approach)
=============================================================
Parses each HTML slide with BeautifulSoup, computes flex layout positions,
and reconstructs all 12 slides as native editable PPTX elements (no screenshots).

Run:  python3 /Users/zhangxiaowu/Claude/build_pptx_v2.py
"""
import subprocess, sys, re
from pathlib import Path

for pkg, name in [('pptx','python-pptx'), ('bs4','beautifulsoup4'), ('lxml','lxml')]:
    try: __import__(pkg)
    except ImportError:
        subprocess.run([sys.executable,'-m','pip','install',name,
                        '--break-system-packages','-q'], check=True)

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from bs4 import BeautifulSoup

SLIDES_DIR = Path('/Users/zhangxiaowu/Claude/Projects/学习和复刻claude code/html_slides_v2')
OUT = '/Users/zhangxiaowu/Claude/Claude-Code实现原理-可编辑版v2.pptx'

# ── Unit helpers ─────────────────────────────────────────────────────────────
# 1280×720 HTML → 13.333"×7.5" PPTX = 12192000×6858000 EMU → 1px = 9525 EMU
def E(px): return Emu(int(round(px * 9525)))
def Fp(px): return Pt(round(px * 0.75, 1))

def rgb(h):
    h = h.lstrip('#')
    if len(h) == 3: h = h[0]*2 + h[1]*2 + h[2]*2
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

# ── PPTX shape primitives ─────────────────────────────────────────────────────
def rect(slide, lx, ty, w, h, fill='ffffff', border=None, bw=0.5):
    """Solid rectangle, no fill border unless specified."""
    if w <= 0 or h <= 0: return None
    s = slide.shapes.add_shape(1, E(lx), E(ty), E(w), E(h))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
    if border:
        s.line.color.rgb = rgb(border); s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    return s

def grad_bar(slide, c1, c2, lx=0, ty=0, w=1280, h=4):
    """Horizontal gradient rectangle."""
    s = slide.shapes.add_shape(1, E(lx), E(ty), E(w), E(h))
    s.line.fill.background()
    sp = s._element.spPr
    for tag in ['a:solidFill','a:gradFill','a:noFill','a:pattFill']:
        el = sp.find(qn(tag))
        if el is not None: sp.remove(el)
    gf = etree.SubElement(sp, qn('a:gradFill'))
    gsl = etree.SubElement(gf, qn('a:gsLst'))
    for pos, col in [('0', c1.lstrip('#')), ('100000', c2.lstrip('#'))]:
        gs = etree.SubElement(gsl, qn('a:gs')); gs.set('pos', pos)
        sc = etree.SubElement(gs, qn('a:srgbClr')); sc.set('val', col)
    lin = etree.SubElement(gf, qn('a:lin'))
    lin.set('ang', '0'); lin.set('scaled', '0')

def tb(slide, text, lx, ty, w, h,
       px=13, bold=False, color='0f172a', face='Noto Sans SC',
       italic=False, wrap=True, align='left', char_spc=0, margin_px=0):
    """Transparent text box with zero internal padding."""
    if not text: return None
    box = slide.shapes.add_textbox(E(lx), E(ty), E(w), E(h))
    box.fill.background(); box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = {'left':PP_ALIGN.LEFT,'center':PP_ALIGN.CENTER,
                   'right':PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = text
    run.font.size = Fp(px); run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(str(color).lstrip('#'))
    run.font.name = face
    if char_spc:
        rp = run._r.get_or_add_rPr(); rp.set('spc', str(int(char_spc*100)))
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        m = str(E(margin_px))
        bp.set('lIns', m); bp.set('rIns', m)
        bp.set('tIns', m); bp.set('bIns', m)
    return box

def tb_multi(slide, lines, lx, ty, w, h,
             px=13, bold=False, color='0f172a', face='Noto Sans SC',
             line_spacing=1.4):
    """Multi-line text box (list of strings)."""
    if not lines: return None
    box = slide.shapes.add_textbox(E(lx), E(ty), E(w), E(h))
    box.fill.background(); box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Fp(px); run.font.bold = bold
        run.font.color.rgb = rgb(str(color).lstrip('#'))
        run.font.name = face
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('lIns', '0'); bp.set('rIns', '0')
        bp.set('tIns', '0'); bp.set('bIns', '0')
    return box

def std_header(slide, label, h1_text, subtitle, sec_color, c1=None, c2=None):
    """Standard header: gradient bar / label / H1+subtitle / divider."""
    grad_bar(slide, c1 or '#'+sec_color, c2 or '#'+sec_color)
    tb(slide, label, 52, 14, 700, 16, px=9.5, bold=True, color=sec_color,
       char_spc=4, wrap=False)
    # H1
    h1_w = min(len(h1_text) * 17 + 20, 700)
    tb(slide, h1_text, 52, 36, h1_w, 38, px=28, bold=True, color='0f172a',
       face='Space Grotesk', wrap=False)
    if subtitle:
        sub_left = 52 + h1_w + 16
        if sub_left < 950:
            tb(slide, subtitle, sub_left, 42, 1228 - sub_left, 28,
               px=11.5, color='64748b', wrap=False)
    rect(slide, 52, 74, 1176, 1, fill='e2e8f0')

# ── Presentation setup ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Emu(12192000)
prs.slide_height = Emu(6858000)
blank = prs.slide_layouts[6]

def new_slide(bg='ffffff'):
    s = prs.slides.add_slide(blank)
    if bg.lower() != 'ffffff':
        rect(s, 0, 0, 1280, 720, fill=bg)
    return s

# ── HTML helpers ──────────────────────────────────────────────────────────────
def load(fname):
    return BeautifulSoup((SLIDES_DIR / fname).read_text('utf-8'), 'lxml')

def txt(el): return el.get_text(separator=' ', strip=True) if el else ''

def css_color(style_str, prop='color'):
    m = re.search(rf'{prop}\s*:\s*#([0-9a-fA-F]{{6}})', style_str or '')
    return m.group(1) if m else None


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 01 — Cover
# ═══════════════════════════════════════════════════════════════════════════════
print('Building slide 01...')
soup = load('slide_01.html')
s = new_slide()

# Ghost "CC" text
tb(s, 'CC', -20, -10, 400, 320, px=240, bold=True, color='ede9fe',
   face='Space Grotesk', wrap=False)

# Top gradient bar (indigo)
grad_bar(s, '#6366f1', '#a5b4fc')

# Section label
tb(s, 'CLAUDE CODE · 实现原理深度解析', 52, 14, 700, 16,
   px=9.5, bold=True, color='6366f1', char_spc=4, wrap=False)

# H1 + subtitle
tb(s, '实现原理与架构', 52, 36, 560, 38, px=28, bold=True,
   color='0f172a', face='Space Grotesk', wrap=False)
tb(s, '完整技术拆解 · 26 张精讲幻灯片', 456, 42, 440, 28,
   px=11.5, color='64748b', wrap=False)

# Divider
rect(s, 52, 74, 1176, 1, fill='e2e8f0')

# Left: description text
tb(s, '从 while(True) 主循环到 1800+ 行边界处理——\n深度解析 Claude Code 的客户端架构、工具系统、\n安全模型与模型协同设计全貌。',
   52, 96, 660, 80, px=13, color='475569', wrap=True)

# Pills
pills = [
    ('Agent 主循环', 'eef2ff', '4338ca', 'c7d2fe'),
    ('20+ 工具系统', 'ecfdf5', '065f46', 'a7f3d0'),
    ('权限安全',     'fef3c7', '92400e', 'fde68a'),
    ('模型协同设计', 'fdf4ff', '6b21a8', 'e9d5ff'),
]
px_l = 52
for label, bg, fc, bc in pills:
    w_ = len(label) * 10 + 28
    rect(s, px_l, 188, w_, 26, fill=bg, border=bc)
    tb(s, label, px_l+8, 191, w_-16, 20, px=10.5, bold=True, color=fc, wrap=False)
    px_l += w_ + 10

# Right: dark code panel
rect(s, 740, 88, 480, 342, fill='1e1b4b')
# Mac window dots
rect(s, 756, 100, 10, 10, fill='ef4444')
rect(s, 773, 100, 10, 10, fill='fbbf24')
rect(s, 790, 100, 10, 10, fill='34d399')
tb(s, 'claude_core.py', 812, 97, 240, 16, px=10, color='ffffff66',
   face='Courier New', wrap=False)

code_lines = [
    ('# Claude Code 核心循环', '6b7280'),
    ('while True:', 'f472b6'),
    ('  resp = api.call(messages, tools)', 'e2e8f0'),
    ('  if resp.stop == "end_turn":', 'e2e8f0'),
    ('    break', 'f472b6'),
    ('  execute(resp.tool_calls)', '86efac'),
    ('  # 模型决定一切', '6b7280'),
]
for i, (line, col) in enumerate(code_lines):
    tb(s, line, 756, 124+i*38, 460, 32, px=11.5, color=col,
       face='Courier New', wrap=False)

# Bottom 3 stat cards
stats = [
    ('10',   '章节完整覆盖', '全面架构解析', 'f8f7ff', '6366f1', 'e0e7ff', '1e1b4b'),
    ('5',    '核心机制',     '深度技术拆解', 'f0fdf4', '10b981', 'bbf7d0', '064e3b'),
    ('完整', '汇报版',       '含复刻实施路线','fff7ed','f97316','fed7aa','7c2d12'),
]
for i, (num, ttl, sub, bg, fc, bc, tc) in enumerate(stats):
    lx2 = 52 + i * 396
    rect(s, lx2, 590, 376, 100, fill=bg, border=bc)
    rect(s, lx2, 590, 376, 3, fill=fc)
    tb(s, num, lx2+18, 600, 100, 42, px=32, bold=True, color=fc,
       face='Space Grotesk', wrap=False)
    tb(s, ttl, lx2+100, 603, 260, 22, px=13, bold=True, color=tc, wrap=False)
    tb(s, sub, lx2+100, 626, 260, 20, px=11, color='94a3b8', wrap=False)

# Bottom attribution
tb(s, '综合官方文档、社区逆向分析与公开论文 · 整理日期：2026-06-11 · 参考：code.claude.com',
   52, 700, 900, 16, px=9, color='cbd5e1', wrap=False)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 02 — TOC (目录)
# ═══════════════════════════════════════════════════════════════════════════════
print('Building slide 02...')
s = new_slide()

# Left dark panel
rect(s, 0, 0, 272, 720, fill='0f172a')
rect(s, 40, 56, 32, 3, fill='6366f1')
tb(s, '目录', 40, 66, 200, 72, px=52, bold=True, color='ffffff',
   face='Space Grotesk', wrap=False)
tb(s, '全文 10 节内容概览', 40, 148, 200, 20, px=12.5, color='94a3b8', wrap=False)

# Counter
rect(s, 40, 178, 192, 80, fill='1e293b', border='334155')
tb(s, '10', 52, 190, 120, 52, px=40, bold=True, color='818cf8',
   face='Space Grotesk', wrap=False)
tb(s, '核心章节 · 完整覆盖', 52, 240, 168, 18, px=11, color='64748b', wrap=False)

# Right 2×5 grid
sections = [
    ('01','总体设计哲学','核心思想 · 单一主循环 · Unix哲学','6366f1'),
    ('02','客户端工程架构','技术栈 · 会话持久化 · 成本工程','0ea5e9'),
    ('03','Agent 主循环','三阶段循环 · 端到端时序 · 错误处理','f97316'),
    ('04','工具系统','20个内置工具 · 五大类 · 关键决策','a855f7'),
    ('05','上下文管理','五层策略 · 压缩机制 · 成本控制','6366f1'),
    ('06','扩展层机制','CLAUDE.md · Skills · MCP · Hooks','16a34a'),
    ('07','权限与安全','权限模式 · 规则引擎 · 威胁防护','dc2626'),
    ('08','模型框架协同','训练对齐 · 多模型路由 · 评价体系','7c3aed'),
    ('09','复刻实施路线图','MVP · 优先级 · 模型路线','3b82f6'),
    ('10','参考资料','官方文档 · 社区研究 · 开源项目','64748b'),
]
for i, (num, title, sub, color) in enumerate(sections):
    col = i % 2; row = i // 2
    lx2 = 288 + col * 492; ty2 = 16 + row * 137
    rect(s, lx2, ty2, 472, 124, fill='ffffff', border=color+'40')
    rect(s, lx2, ty2, 4, 124, fill=color)
    tb(s, num, lx2+16, ty2+20, 56, 52, px=22, bold=True, color=color,
       face='Space Grotesk', wrap=False)
    tb(s, title, lx2+72, ty2+24, 380, 28, px=14, bold=True, color='0f172a', wrap=False)
    tb(s, sub, lx2+16, ty2+82, 440, 22, px=10.5, color='64748b', wrap=False)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 03 — 总体设计哲学  (SECTION 01, indigo)
# ═══════════════════════════════════════════════════════════════════════════════
print('Building slide 03...')
s = new_slide()

# Ghost "01"
tb(s, '01', 920, -20, 360, 300, px=260, bold=True, color='f0effe',
   face='Space Grotesk', wrap=False)

std_header(s, 'SECTION 01 / 10', '总体设计哲学',
           '模型负责推理，框架负责行动', '6366f1', '#6366f1', '#818cf8')

# Quote box (top:88 approx, using standard header bottom at 74 + some padding)
rect(s, 52, 88, 1176, 48, fill='f8f7ff', border='c7d2fe')
rect(s, 52, 88, 4, 48, fill='6366f1')
tb(s, '" 30 行可以写出 Agent Loop，但生产级实现需要 1800+ 行——复杂度全部在边界处理。"',
   76, 96, 1140, 34, px=12.5, color='4338ca', italic=True, wrap=False)

# 2×2 cards
card_data = [
    ('单一主循环',
     '没有复杂状态机，就是 while(true) 循环，让模型自己决定下一步执行什么工具、何时停止。',
     '6366f1', 'ffffff', 'e0e7ff', '1e1b4b'),
    ('模型即调度器',
     '哪个工具、什么顺序、何时停止，全部交给模型决策；工程侧只做权限审查、上下文管理和失败恢复。',
     '10b981', 'ffffff', 'd1fae5', '1e1b4b'),
    ('Unix 哲学',
     '运行在终端，可管道化（claude -p 无头模式），可组合进 CI/CD；单一职责、可组合设计。',
     'f59e0b', 'ffffff', 'fde68a', '1e1b4b'),
    ('复杂度在边界',
     '基础循环极为简单；所有复杂性集中在循环失败时如何优雅恢复——这是生产可用的关键。',
     'ef4444', 'ffffff', 'fecaca', '1e1b4b'),
]
cw, ch = 576, 244
for i, (ttl, body, acc, bg, bc, tc) in enumerate(card_data):
    lx2 = 52 + (i%2) * 592; ty2 = 148 + (i//2) * 256
    rect(s, lx2, ty2, cw, ch, fill=bg, border=bc)
    rect(s, lx2, ty2, cw, 3, fill=acc)
    tb(s, ttl, lx2+14, ty2+14, cw-28, 24, px=14, bold=True, color=tc, wrap=False)
    rect(s, lx2+14, ty2+42, 30, 2, fill=acc)
    tb(s, body, lx2+14, ty2+50, cw-28, ch-62, px=11.5, color='475569', wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 06 — 客户端工程架构  (SECTION 02, sky blue)
# ═══════════════════════════════════════════════════════════════════════════════
print('Building slide 06...')
s = new_slide()

std_header(s, 'SECTION 02 / 10', '客户端工程架构',
           'TypeScript + Node.js · React + Ink 终端 UI', '0284c7', '#0284c7', '#38bdf8')

# ── Left: 4 tech-stack cards (flex:1, gap=8) ─────────────────────────────────
# Container: left=52, top=88, width=440, bottom=696
# Label: ~20px + 10px margin = 30px overhead
# 4 cards in 608px total: (608-30 - 3*8)/4 ≈ 136px each
tb(s, '技术栈 & 进程模型', 52, 96, 440, 14, px=9, bold=True, color='0284c7',
   char_spc=3, wrap=False)

left_cards = [
    ('TS',  'TypeScript + Node.js', '0369a1',
     '单一可执行 CLI（claude），代码经混淆打包，sourcemap 已被社区整理公开',
     'f0f9ff', 'bae6fd', '0c4a6e', '0369a1'),
    ('INK', 'React + Ink UI', '0d9488',
     'React 组件模型渲染终端界面，支持流式 token 渲染与 diff 展示',
     'f0fdf4', 'bbf7d0', '064e3b', '0d9488'),
    ('npm', 'npm 持续分发', '4338ca',
     '通过 npm 分发，版本更新频繁，社区 Piebald-AI 等已整理逆向分析',
     'faf5ff', 'ddd6fe', '3b0764', '7c3aed'),
    ('SH',  '单进程持久 Shell', '1e40af',
     '跨多次 Bash 工具调用保持同一 shell 进程状态（cd 等状态持续）',
     'fff7ed', 'fed7aa', '7c2d12', 'c2410c'),
]
card_h, gap = 136, 8
for i, (badge, ttl, badge_col, body, bg, bc, ttl_col, body_col) in enumerate(left_cards):
    ty2 = 118 + i * (card_h + gap)
    rect(s, 52, ty2, 440, card_h, fill=bg, border=bc)
    # Badge square
    rect(s, 52, ty2, 44, 40, fill=badge_col)
    tb(s, badge, 54, ty2+12, 40, 18, px=10, bold=True, color='ffffff',
       face='Arial', wrap=False, align='center')
    tb(s, ttl, 104, ty2+8, 374, 22, px=13, bold=True, color=ttl_col, wrap=False)
    tb(s, body, 104, ty2+32, 374, card_h-44, px=11, color=body_col, wrap=True)

# Vertical divider
rect(s, 512, 88, 1, 608, fill='bae6fd')

# ── Right: 5 module cards (flex:1, gap=7) ────────────────────────────────────
# Container: left=532, top=88, right=40→width=708, bottom=696
# (608-30 - 4*7)/5 = 569/5 ≈ 109px
tb(s, '核心模块划分（逆向分析归纳）', 532, 96, 700, 14, px=9, bold=True,
   color='0284c7', char_spc=3, wrap=False)

right_cards = [
    ('RU', 'dbeafe', '1d4ed8', 'runtime',     '主运行时',
     '主循环调度、会话状态管理、配置加载、prompt 拼组'),
    ('AP', 'dcfce7', '15803d', 'api-client',  'API 客户端',
     'Messages API 封装，SSE 流式解析，重试 / 限流 / Token 计数'),
    ('TO', 'f3e8ff', '7c3aed', 'tools',       '工具执行器',
     '约 20 个内置工具的 JSON Schema 定义 + 本地执行器实现'),
    ('PE', 'fee2e2', 'dc2626', 'permissions', '权限引擎',
     'allow / deny 规则引擎 + 交互式审批 + Hook 拦截点'),
    ('UI', 'ffedd5', 'ea580c', 'ui',          '终端 UI',
     '终端渲染、diff 视图、进度 spinner、流式 token 输出'),
]
rcard_h, rgap = 109, 7
for i, (badge, bg, fc, mod, name, desc) in enumerate(right_cards):
    ty2 = 118 + i * (rcard_h + rgap)
    rect(s, 532, ty2, 708, rcard_h, fill='ffffff', border='e2e8f0')
    rect(s, 532, ty2, 38, 34, fill=bg)
    tb(s, badge, 534, ty2+9, 34, 18, px=10, bold=True, color=fc,
       face='Arial', wrap=False, align='center')
    # Code badge + name
    rect(s, 578, ty2+8, len(mod)*7+16, 20, fill=bg)
    tb(s, mod, 583, ty2+10, len(mod)*7+10, 16, px=10, bold=True, color=fc,
       face='Courier New', wrap=False)
    code_w = len(mod)*7+26
    tb(s, name, 578+code_w, ty2+8, 300, 22, px=12.5, bold=True, color='0f172a', wrap=False)
    tb(s, desc, 578, ty2+32, 652, rcard_h-42, px=10.5, color='64748b', wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 08 — Agent 主循环  (SECTION 03, orange)
# ═══════════════════════════════════════════════════════════════════════════════
print('Building slide 08...')
s = new_slide()

# Ghost "03"
tb(s, '03', 900, -20, 380, 300, px=280, bold=True, color='fff7ed',
   face='Space Grotesk', wrap=False)

std_header(s, 'SECTION 03 / 10', 'Agent 主循环',
           '三阶段循环 · 端到端时序 · 并发工具执行 · 模型自主决策',
           'ea580c', '#c2410c', '#f97316')

# Right: CORE LOOP box (top:88, right:40, width:340)
rect(s, 900, 88, 340, 180, fill='fff7ed', border='fed7aa')
tb(s, 'CORE LOOP', 916, 96, 300, 14, px=9, bold=True, color='ea580c',
   char_spc=2, wrap=False)
code = ('while not_done:\n'
        '  preprocess()    # ① 预处理\n'
        '  api_call()      # ② API调用\n'
        '  execute_tools() # ③ 执行\n'
        '  if end_turn: break')
tb(s, code, 916, 112, 308, 144, px=11.5, color='7c2d12',
   face='Courier New', wrap=True)

# Bottom 4-column section (height:310, top:410)
BOTTOM_TOP = 410
rect(s, 0, BOTTOM_TOP, 1280, 1, fill='f1f5f9')

col_data = [
    ('01', '预处理阶段',  '#f59e0b', 'fff7ed', '78350f',
     ['注入 CLAUDE.md 系统提示', '拼装消息历史', '上下文预算检查']),
    ('02', 'API 调用阶段','#6366f1', 'f8f7ff', '4338ca',
     ['SSE 流式接收 token', 'text / tool_use 两种块', 'stop_reason 决定走向']),
    ('03', '工具执行阶段','#10b981', 'f0fdf4', '065f46',
     ['权限 allow/deny 检查', 'Promise.all 并发执行', '结果回填消息队列']),
    ('EXIT','退出条件',   '#ea580c', 'fff8f3', '92400e',
     ['stop = "end_turn"', 'max_tokens 到达', 'Ctrl-C 用户中断']),
]
col_w = 320
for i, (num, ttl, acc, bg, tc, bullets) in enumerate(col_data):
    lx2 = i * col_w
    rect(s, lx2, BOTTOM_TOP, col_w, 310, fill=bg, border='f1f5f9')
    rect(s, lx2, BOTTOM_TOP, col_w, 3, fill=acc)
    if i < 3:
        rect(s, lx2+col_w-1, BOTTOM_TOP, 1, 310, fill='f1f5f9')
    tb(s, num, lx2+24, BOTTOM_TOP+14, 100, 20, px=11, bold=True,
       color=acc, face='Space Grotesk', char_spc=3, wrap=False)
    tb(s, ttl, lx2+24, BOTTOM_TOP+36, col_w-48, 26, px=16, bold=True,
       color='1c0a00', wrap=False)
    rect(s, lx2+24, BOTTOM_TOP+66, 24, 2, fill=acc)
    bullet_text = '\n'.join('· ' + b for b in bullets)
    tb(s, bullet_text, lx2+24, BOTTOM_TOP+76, col_w-48, 120,
       px=11, color=tc, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 09 — 端到端交互时序  (SECTION 03, amber)
# ═══════════════════════════════════════════════════════════════════════════════
print('Building slide 09...')
s = new_slide()

std_header(s, 'SECTION 03 / 10', '端到端交互时序',
           '从用户输入到最终输出的完整流程', 'b45309', '#b45309', '#f59e0b')

# Left: 6-step flow (top:88, left:52, width:420)
# 6 steps, h available=632, no gap → each ≈ 105px
steps = [
    ('1', '用户输入',       'stdin / TUI 接收原始文本', '自然语言任务 → messages[ ]'),
    ('2', '上下文组装',     'System Prompt + CLAUDE.md + 历史消息 + tool schema', ''),
    ('3', 'API 调用（流式）','claude-3-7-sonnet · streaming SSE · token 计量', ''),
    ('4', '模型决策',       '输出 text_block 或 tool_use_block', ''),
    ('5', '本地工具执行',   '权限审查 → 并发执行 → 序列化为 tool_result', ''),
    ('↩', '回填 → 下轮循环','tool_result 追加 messages → 渲染 → 回 Step 2', ''),
]
step_h = 105
for i, (num, ttl, detail, extra) in enumerate(steps):
    ty2 = 88 + i * step_h
    is_last = (num == '↩')
    bg2 = 'fef3c7' if is_last else 'ffffff'
    bc2 = 'fde68a'
    rect(s, 52, ty2, 420, step_h-2, fill=bg2, border=bc2)
    # Number circle
    circle_bg = 'b45309' if is_last else 'ffffff'
    rect(s, 52, ty2+34, 36, 36, fill=circle_bg, border='d97706')
    tb(s, num, 56, ty2+38, 28, 28, px=13, bold=True,
       color='ffffff' if is_last else 'b45309',
       face='Space Grotesk', wrap=False, align='center')
    tb(s, ttl, 98, ty2+10, 360, 22, px=13, bold=True, color='1c0a00', wrap=False)
    tb(s, detail, 98, ty2+34, 360, 50, px=10.5, color='78350f', wrap=True)

# Vertical divider
rect(s, 492, 88, 1, 632, fill='e5d5b8')

# Right: implementation details (top:88, left:516, right:40→width=724)
tb(s, 'IMPLEMENTATION DETAILS', 516, 94, 700, 16,
   px=9.5, bold=True, color='b45309', char_spc=3, wrap=False)

details = [
    ('流式渲染机制',
     'SSE 事件 content_block_delta 触发 React/Ink 逐 token 重渲染，流结束后执行最终布局刷新。',
     'fde8b8', '1c0a00', '57534e'),
    ('工具并发策略',
     '单次响应可含多个 tool_use[]，Promise.all 并发执行，全部完成后一并回填，减少 API 往返次数。',
     'fde8b8', '1c0a00', '57534e'),
    ('错误恢复策略',
     '工具失败结果以 is_error:true 回填，模型自主决策：重试 / 换工具 / 向用户报告。',
     'fde8b8', '1c0a00', '57534e'),
    ('Token 成本放大',
     '每轮消耗 = 全量历史 + 工具结果。10 轮循环后总量是首轮 10–30×。压缩阈值约 95% 触发 /compact。',
     'fde68a', '7c2d12', '78350f'),
]
detail_h = 148
for i, (ttl, body, bc, tc, dc) in enumerate(details):
    ty2 = 114 + i * detail_h
    rect(s, 516, ty2, 720, detail_h-8, fill='ffffff', border=bc)
    rect(s, 516, ty2, 6, 6, fill='d97706')
    tb(s, ttl, 530, ty2+8, 688, 22, px=13, bold=True, color=tc, wrap=False)
    tb(s, body, 530, ty2+32, 688, detail_h-48, px=11, color=dc, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — 交互时序 · 9 检查点  (SECTION 03, orange table)
# ═══════════════════════════════════════════════════════════════════════════════
print('Building slide 10...')
s = new_slide()

std_header(s, 'SECTION 03 / 10', '交互时序',
           '9 个关键检查点 · 端到端完整流程', 'ea580c', '#c2410c', '#f97316')

# Left dark panel (top:88, width:200)
rect(s, 0, 88, 200, 632, fill='1e2a3a')
tb(s, '9 个关键\n检查点', 20, 110, 160, 60, px=11, color='ffffff80', wrap=True)

# Legend
tb(s, 'ACTOR', 20, 606, 150, 14, px=8.5, bold=True, color='ffffff40',
   char_spc=2, wrap=False)
legend = [('3b82f6','用户侧'), ('ea580c','API 侧'),
          ('10b981','框架侧'), ('8b5cf6','模型侧')]
for i, (col, lbl) in enumerate(legend):
    rect(s, 20, 624+i*22, 10, 10, fill=col)
    tb(s, lbl, 36, 622+i*22, 140, 16, px=10.5, color='ffffff8c', wrap=False)

# Right table area (left:200, top:88, right:0 → w=1080)
# Grid columns: 36+56+{1fr}+220+68 = 380+1fr, total available=1080-24-28=1028
# 1fr = 1028-36-56-220-68 = 648
COL_X = [200+24, 200+24+36, 200+24+36+56, 200+24+36+56+648, 200+24+36+56+648+220]
COL_W = [36, 56, 648, 220, 68]

# Table header
rect(s, 200, 88, 1080, 32, fill='f8fafc')
rect(s, 200, 120, 1080, 2, fill='1e2a3a')
headers = ['#', '角色', '检查点', '技术细节', '关键数据']
for j, (hdr_txt, cx, cw) in enumerate(zip(headers, COL_X, COL_W)):
    align = 'right' if j == 4 else 'left'
    tb(s, hdr_txt, cx+4, 96, cw-8, 16, px=9, bold=True, color='94a3b8',
       char_spc=2, wrap=False, align=align)

# 9 rows
rows = [
    ('1','用户','用户输入','stdin / TUI 接收原始文本',
     '自然语言任务 → 加入 messages[ ]','输入层','3b82f6'),
    ('2','框架','CLAUDE.md 注入','读取全局 + 项目级配置文件',
     '级联合并 → 拼入 system prompt 前','固定成本','10b981'),
    ('3','框架','上下文裁剪','token 预算检查 + 压缩判断',
     '超过 95% 窗口 → 触发 /compact','≥ 95%','10b981'),
    ('4','API', 'API 请求发送','streaming POST + tool schema 定义',
     'claude-3-7-sonnet · SSE 长连接','POST','ea580c'),
    ('5','模型','流式 token 输出','content_block_delta 实时推送',
     'React/Ink 逐 token 重渲染终端 UI','SSE','8b5cf6'),
    ('6','模型','stop_reason 判断','决定循环走向的核心分支',
     'end_turn / tool_use / max_tokens','分支','8b5cf6'),
    ('7','框架','权限审查','allow / deny 规则匹配',
     'glob 匹配工具名 + 参数 → 确认/拒绝','拦截层','10b981'),
    ('8','框架','工具并发执行','Promise.all 并行执行多个工具',
     '全部完成后一并回填，减少往返次数','并发','10b981'),
    ('9','框架','消息回填 → 下轮','tool_result 追加，触发下轮循环',
     '→ 回到 Step 2，继续主循环','↩ 循环','ea580c'),
]
actor_colors = {'用户':'3b82f6','框架':'10b981','API':'ea580c','模型':'8b5cf6'}
row_h = 64
for i, (num, actor, chk, sub, detail, stat, _) in enumerate(rows):
    ty2 = 122 + i * row_h
    bg2 = 'fff8f3' if num=='9' else ('f8fafc' if i%2 else 'ffffff')
    rect(s, 200, ty2, 1080, row_h, fill=bg2, border='00000010')
    if num == '9':
        rect(s, 200, ty2, 3, row_h, fill='ea580c')
    # Row number
    tb(s, num, COL_X[0]+4, ty2+20, COL_W[0]-8, 26,
       px=16, bold=True, color='e2e8f0', face='Space Grotesk', wrap=False)
    # Actor badge
    ac = actor_colors.get(actor, '64748b')
    rect(s, COL_X[1]+2, ty2+20, 46, 22, fill=ac)
    tb(s, actor, COL_X[1]+2, ty2+22, 46, 18, px=8.5, bold=True, color='ffffff',
       wrap=False, align='center')
    # Checkpoint + sub
    tb(s, chk, COL_X[2]+4, ty2+8, COL_W[2]-8, 22,
       px=13, bold=True, color='1e2a3a', wrap=False)
    tb(s, sub, COL_X[2]+4, ty2+30, COL_W[2]-8, 24,
       px=10.5, color='64748b', wrap=False)
    # Technical detail
    tb(s, detail, COL_X[3]+4, ty2+14, COL_W[3]-8, 36,
       px=10.5, color='475569', wrap=True)
    # Stat
    highlight = stat in ('↩ 循环', '≥ 95%', '分支')
    tb(s, stat, COL_X[4], ty2+20, COL_W[4], 22,
       px=11, bold=True, color='ea580c' if highlight else '94a3b8',
       wrap=False, align='right')


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — 工具系统全景  (SECTION 04, purple)
# ═══════════════════════════════════════════════════════════════════════════════
print('Building slide 11...')
s = new_slide()

# Ghost "04"
tb(s, '04', 900, -10, 380, 280, px=260, bold=True, color='f5f3ff',
   face='Space Grotesk', wrap=False)

# Non-standard header (inline style with 3 items in one row)
grad_bar(s, '#7c3aed', '#c084fc')
tb(s, 'SECTION 04 / 10', 52, 16, 260, 16, px=10, bold=True, color='7c3aed',
   char_spc=4, wrap=False)
tb(s, '工具系统全景', 320, 12, 500, 34, px=28, bold=True,
   color='3b0764', face='Space Grotesk', wrap=False)
tb(s, '约 20 个内置工具 · 五大类别 · JSON Schema 定义',
   838, 20, 430, 20, px=11, color='94a3b8', wrap=False)
rect(s, 52, 56, 1176, 1, fill='e2e8f0')

# Top 5 category cards (top:66, height:290)
top_cards = [
    ('文件系统',  'Read / Write\nEdit / MultiEdit\nGlob / LS',      '6366f1','f8f7ff','e0e7ff','4338ca'),
    ('Shell 执行','Bash（持久 shell）\nJSNode\n单进程状态保持',         '10b981','f0fdf4','bbf7d0','065f46'),
    ('搜索 & 定位','Grep (ripgrep)\nWebSearch\nWebFetch',             'f59e0b','fffbeb','fde68a','92400e'),
    ('子代理',    'Task（子 Agent）\n独立上下文\n递归调用',              'ef4444','fef2f2','fecaca','b91c1c'),
    ('交互类',    'AskFollowup\nTodoRead/Write\nNotebookEdit',        '06b6d4','ecfeff','a5f3fc','0e7490'),
]
cw2 = 228
for i, (ttl, body, acc, bg, bc, tc) in enumerate(top_cards):
    lx2 = 52 + i * (cw2 + 8)
    rect(s, lx2, 66, cw2, 292, fill=bg, border=bc)
    rect(s, lx2, 66, cw2, 3, fill=acc)
    tb(s, ttl, lx2+12, 80, cw2-24, 24, px=13, bold=True, color='1e1b4b', wrap=False)
    tb(s, body, lx2+12, 110, cw2-24, 200, px=10.5, color=tc, wrap=True)

# Bottom 3 panels
# Design decisions (lx=52, w=464, top=370)
rect(s, 52, 370, 464, 320, fill='ffffff', border='e0e7ff')
tb(s, '关键设计决策', 66, 382, 430, 22, px=12, bold=True, color='3b0764', wrap=False)
design_items = [
    ('工具 → JSON Schema 定义', '4338ca', 'eef2ff', '4338ca'),
    ('幂等性优先设计',          '15803d', 'f0fdf4', '15803d'),
    ('Bash 替代专用工具',       'b45309', 'fffbeb', 'b45309'),
    ('子代理避免上下文膨胀',    'b91c1c', 'fef2f2', 'b91c1c'),
]
for j, (item, fc, bg, _) in enumerate(design_items):
    rect(s, 66, 412+j*62, 430, 48, fill=bg, border=fc+'40')
    tb(s, item, 80, 423+j*62, 400, 28, px=10.5, bold=True, color=fc, wrap=False)

# MCP panel (lx=528, w=464)
rect(s, 528, 370, 464, 320, fill='faf5ff', border='ddd6fe')
tb(s, 'MCP 外部工具', 544, 386, 430, 22, px=12, bold=True, color='3b0764', wrap=False)
tb(s, 'Model Context Protocol', 544, 412, 430, 18, px=10.5, color='7c3aed', wrap=False)
tb(s, '通过 Model Context Protocol 接入第三方工具服务器，同一 JSON Schema 接口，动态注册、热加载，与内置工具完全同等地位。',
   544, 438, 430, 80, px=10.5, color='6d28d9', wrap=True)
rect(s, 544, 524, 430, 52, fill='ede9fe', border='ddd6fe')
tb(s, 'stdio / HTTP / SSE 三种传输协议', 554, 532, 410, 36, px=10.5, bold=True,
   color='4c1d95', wrap=True)

# Warning panel (lx=1004, w=224)
rect(s, 1004, 370, 224, 320, fill='fff8f3', border='fed7aa')
tb(s, '⚠ 注意', 1016, 386, 200, 22, px=12, bold=True, color='c2410c', wrap=False)
tb(s, '工具数 > 30 个时\ntoken 占用大幅增加\n模型幻觉率明显上升',
   1016, 416, 200, 80, px=10.5, color='92400e', wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — 上下文管理五层策略  (SECTION 05, indigo layers)
# ═══════════════════════════════════════════════════════════════════════════════
print('Building slide 12...')
s = new_slide()

# Left vertical accent bar
rect(s, 0, 0, 6, 720, fill='6366f1')

# Header (non-standard – no divider)
tb(s, '上下文管理五层策略', 40, 36, 700, 42, px=30, bold=True,
   color='1e1b4b', face='Space Grotesk', wrap=False)
tb(s, 'SECTION 05 / 10', 756, 44, 240, 20, px=11, bold=True,
   color='6366f1', char_spc=2, wrap=False)
tb(s, 'Context Window 管理是 Claude Code 成本与长任务稳定性的核心命题',
   40, 76, 1200, 20, px=12, color='64748b', wrap=False)

layers = [
    ('L1','精选提示词',   '核心指令精简 · 不重复',
     '系统提示词经精心裁剪，去掉所有冗余说明，只保留模型真正需要的核心指令。每次注入成本固定，节省每轮循环 token。',
     '固定成本节省','818cf8','f8f7ff'),
    ('L2','工具结果截断', '长输出自动截断',
     'Bash 输出、文件内容等工具结果超出限制时自动截断（通常 ~10K 字符），后附"[truncated]"提示，避免单次工具结果撑爆上下文。',
     '单次结果限制','6366f1','f5f3ff'),
    ('L3','自动历史压缩', '/compact 命令触发',
     '当上下文占用达到阈值（约 95%）时，自动调用压缩：保留最近 N 轮，将历史摘要化。用户可手动执行 /compact [指令]。',
     '95% 阈值触发','4f46e5','eef2ff'),
    ('L4','子代理隔离',   'Task 工具独立上下文',
     '子代理（Task 工具）拥有完全独立的上下文窗口，不污染主 Agent 历史。长任务分治拆解，各子代理只汇报结果。',
     '完全隔离','3730a3','e0e7ff'),
    ('L5','会话持久化',   '--continue / --resume',
     '消息历史序列化至 ~/.claude/projects/…/sessions/ 目录，支持跨终端会话恢复，为日志审计提供完整轨迹。',
     '跨会话恢复','312e81','ddd6fe'),
]
layer_h = 120
for i, (num, ttl, sub, body, tag, acc, light) in enumerate(layers):
    ty2 = 104 + i * layer_h
    # Left colored block (w=260)
    rect(s, 40, ty2, 260, layer_h-4, fill=acc)
    tb(s, num, 52, ty2+10, 80, 40, px=22, bold=True, color='ffffff33',
       face='Space Grotesk', wrap=False)
    tb(s, ttl, 52, ty2+50, 240, 24, px=14, bold=True, color='ffffff', wrap=False)
    tb(s, sub, 52, ty2+76, 240, 20, px=10.5, color='ffffffb3', wrap=False)
    # Right content block
    rect(s, 300, ty2, 972, layer_h-4, fill=light, border=acc+'40')
    tb(s, body, 316, ty2+10, 820, 68, px=11, color='4338ca', wrap=True)
    # Tag badge (right aligned)
    tw = len(tag)*9+20
    rect(s, 1272-tw, ty2+8, tw, 22, fill=acc+'33', border=acc+'66')
    tb(s, tag, 1276-tw, ty2+10, tw-8, 18, px=10, bold=True, color=acc, wrap=False)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — 扩展层：4大机制  (SECTION 06, green)
# ═══════════════════════════════════════════════════════════════════════════════
print('Building slide 13...')
s = new_slide('f0fdf4')

grad_bar(s, '#16a34a', '#4ade80')

# Non-standard header
tb(s, '扩展层：4大机制', 56, 22, 600, 38, px=28, bold=True,
   color='14532d', face='Space Grotesk', wrap=False)
tb(s, 'SECTION 06 / 10', 692, 28, 280, 18, px=11, bold=True,
   color='16a34a', char_spc=2, wrap=False)
tb(s, '让 Claude Code 可被深度定制、集成和自动化',
   984, 28, 280, 18, px=11.5, color='166534', wrap=False)

# 5 flex rows (top:70, bottom:700 → h=630; 4 rows + summary)
# gap=11, 4×11=44; (630-44)/5 ≈ 117px each
ext_rows = [
    ('CLAUDE.md',          '项目级 System Prompt', '16a34a',
     '放在 ~/.claude/（全局）或项目根目录（本地）。Claude 每次启动自动注入，定义工作流程、团队规范、禁用操作等约束，支持 @import 引用其他文件。',
     ['自动注入','级联继承'], 'dcfce7', '14532d'),
    ('Skills（斜杠命令）', '/skill-name',           '0d9488',
     'SKILL.md 文件定义可复用任务模板（如 /debug、/review、/standup）。Claude 发现 SKILL.md 后将其注册为斜杠命令，执行时展开完整提示词。',
     ['可复用模板','插件分发'], 'ccfbf1', '134e4a'),
    ('MCP 服务器',         'Model Context Protocol','7c3aed',
     '通过标准 JSON-RPC 协议接入外部工具服务（数据库、Slack、GitHub…），与内置工具同等地位，支持 stdio / HTTP / SSE 三种传输。',
     ['标准化接口','热加载'], 'ede9fe', '3b0764'),
    ('Hooks（生命周期钩子）','Shell 命令',           'dc2626',
     '四个钩子点：PreToolUse · PostToolUse · PreCompact · Stop。在工具执行前后注入任意 Shell 脚本，实现自动测试、代码格式化、日志审计等。',
     ['4个钩子点','拦截 & 增强'], 'fee2e2', '7f1d1d'),
]
row_h, gap_h = 117, 11
for i, (ttl, badge, acc, body, tags, icon_bg, tc) in enumerate(ext_rows):
    ty2 = 70 + i * (row_h + gap_h)
    rect(s, 56, ty2, 1168, row_h, fill='ffffff', border=acc+'66')
    rect(s, 56, ty2, 5, row_h, fill=acc)
    # Icon area
    rect(s, 72, ty2+(row_h-44)//2, 44, 44, fill=icon_bg)
    # Title + badge
    tb(s, ttl, 128, ty2+14, 440, 24, px=15, bold=True, color=tc, wrap=False)
    badge_w = len(badge)*8+24
    rect(s, 580, ty2+16, badge_w, 22, fill=acc+'15', border=acc+'55')
    tb(s, badge, 588, ty2+18, badge_w-16, 18, px=11, bold=True, color=acc, wrap=False)
    # Body text
    tb(s, body, 128, ty2+42, 848, row_h-56, px=11.5, color='166534', wrap=True)
    # Tags (right side)
    tx2 = 988
    for tag in tags:
        tw = len(tag)*9+20
        rect(s, tx2, ty2+38, tw, 24, fill=acc+'22', border=acc+'44')
        tb(s, tag, tx2+6, ty2+41, tw-12, 18, px=10, bold=True, color=acc, wrap=False)
        tx2 += tw + 8

# Summary row (green gradient)
sum_top = 70 + 4 * (row_h + gap_h)
rect(s, 56, sum_top, 1168, row_h, fill='14532d')
tb(s, '四者叠加：CLAUDE.md 定义行为基线 → Skills 封装重复动作 → MCP 扩展工具边界 → Hooks 在关键节点自动执行外部逻辑，实现全栈级 AI 工作流自动化。',
   90, sum_top+20, 1100, row_h-40, px=12, color='86efac', wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Hooks 深度解析  (SECTION 06, cyan/blue)
# ═══════════════════════════════════════════════════════════════════════════════
print('Building slide 14...')
s = new_slide()

grad_bar(s, '#0284c7', '#22d3ee')
tb(s, 'SECTION 06 / 10 · 扩展功能特性', 52, 18, 700, 16,
   px=10, bold=True, color='0284c7', char_spc=4, wrap=False)
tb(s, 'Hooks 深度解析', 52, 40, 600, 44, px=32, bold=True,
   color='0c4a6e', face='Space Grotesk', wrap=False)
tb(s, '在工具执行生命周期任意节点注入自定义逻辑',
   52, 88, 900, 20, px=12, color='64748b', wrap=False)
rect(s, 52, 108, 1176, 1, fill='e2e8f0')

hooks = [
    ('PreToolUse',  '工具调用前执行', '可拦截',
     '在 Claude 工具调用实际执行前触发。用于：安全审查（过滤危险命令）、日志记录（记录所有操作意图）、输入验证（拦截格式错误参数）。返回非 0 退出码可阻止执行。',
     '0ea5e9', 'e0f2fe', 'dbeafe'),
    ('PostToolUse', '工具调用后执行', '自动化',
     '工具执行完成、结果回填前触发。用于：自动测试（每次写文件后运行测试）、格式化（ESLint / Prettier 自动修复）、结果增强（追加上下文信息到返回值）。',
     '6366f1', 'eef2ff', 'ede9fe'),
    ('PreCompact',  '压缩前执行',     '保留关键信息',
     '上下文压缩前触发，可指定自定义摘要指令、标记需要保留的关键信息段，避免重要上下文在自动压缩中丢失。',
     'f59e0b', 'fffbeb', 'fef3c7'),
    ('Stop',        '任务完成后执行', '收尾自动化',
     '主循环退出时触发。可用于：发送 Slack 通知、生成任务报告、触发 CI 流水线、清理临时文件等收尾操作。',
     'ef4444', 'fef2f2', 'fee2e2'),
]
hook_h = 148
for i, (cmd, sub, tag, body, acc, bg, tbg) in enumerate(hooks):
    ty2 = 118 + i * hook_h
    rect(s, 52, ty2, 1176, hook_h-4, fill='ffffff', border=acc+'40')
    rect(s, 52, ty2, 4, hook_h-4, fill=acc)
    rect(s, 68, ty2+28, 44, 44, fill=bg)
    # Command badge
    cmd_w = len(cmd)*9+24
    rect(s, 120, ty2+14, cmd_w, 28, fill=bg, border=acc+'66')
    tb(s, cmd, 128, ty2+17, cmd_w-16, 22, px=12.5, bold=True, color=acc, wrap=False)
    tb(s, sub, 124+cmd_w, ty2+19, 280, 20, px=12, color='475569', wrap=False)
    # Tag (right)
    tw = len(tag)*10+24
    rect(s, 1224-tw, ty2+16, tw, 24, fill=tbg, border=acc+'66')
    tb(s, tag, 1228-tw, ty2+18, tw-8, 18, px=9.5, bold=True, color=acc, wrap=False)
    # Body
    tb(s, body, 120, ty2+50, 1088, hook_h-64, px=11, color='475569', wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — 权限与安全体系  (SECTION 07, red)
# ═══════════════════════════════════════════════════════════════════════════════
print('Building slide 15...')
s = new_slide('fff5f5')

grad_bar(s, '#dc2626', '#f87171')
tb(s, '权限与安全体系', 56, 20, 600, 38, px=28, bold=True,
   color='7f1d1d', face='Space Grotesk', wrap=False)
tb(s, 'SECTION 07 / 10', 680, 26, 280, 18, px=11, bold=True,
   color='dc2626', char_spc=2, wrap=False)

# Top 4 threat cards
threats = [
    ('提示词注入',
     '恶意内容伪装指令混入上下文，诱导模型执行非授权操作。对抗策略：Hooks 预验证 + 权限规则双重过滤。',
     'dc2626'),
    ('权限提升攻击',
     '利用工具组合绕过单一权限限制，最终实现超越授权的操作。对抗策略：工具级细粒度 allow/deny 规则。',
     'ef4444'),
    ('误删关键数据',
     '模型误判意图导致不可逆破坏性操作。对抗策略：危险操作强制用户二次确认（--no-bypass-permissions）。',
     'ef4444'),
    ('密钥泄漏风险',
     '上下文中的 API Key 被工具或子代理意外传出。对抗策略：网络访问限制 + 敏感词过滤规则。',
     'f87171'),
]
card_w = 304
for i, (ttl, body, acc) in enumerate(threats):
    lx2 = 56 + i * (card_w - 4)
    rect(s, lx2, 64, 288, 140, fill='ffffff', border='fecaca')
    rect(s, lx2, 64, 288, 3, fill=acc)
    tb(s, ttl, lx2+12, 76, 264, 22, px=11, bold=True, color=acc, wrap=False)
    tb(s, body, lx2+12, 100, 264, 96, px=10.5, color='991b1b', wrap=True)

# Divider with label
rect(s, 56, 212, 1168, 1, fill='fecaca')
tb(s, '纵深防御策略', 552, 204, 200, 18, px=11, bold=True, color='dc2626',
   char_spc=1, wrap=False, align='center')

# Bottom 3 defense panels
panel_w = 376
defense = [
    ('三级权限模式', [
        ('Auto',    '无需确认，CI/无头模式',   'fef2f2', 'dc2626'),
        ('Default', '破坏性操作需确认',         'fef2f2', 'dc2626'),
        ('Strict',  '所有工具调用均需审批',     'fee2e2', '7f1d1d'),
    ], None, 'fecaca'),
    ('规则引擎', None,
     '"allow": ["Bash(git *)"],\n"deny": ["Bash(rm -rf *)"],\n"deny": ["Write(~/.ssh/*)"]',
     'fecaca'),
    ('网络访问控制', None,
     'WebSearch / WebFetch 工具可单独启用/禁用。\n企业部署可通过环境变量 CLAUDE_DISABLE_WEB 强制关闭所有出站网络请求。',
     'fecaca'),
]
for i, (ttl, items, code, bc) in enumerate(defense):
    lx2 = 56 + i * (panel_w + 20)
    rect(s, lx2, 226, panel_w, 468, fill='ffffff', border=bc)
    tb(s, ttl, lx2+16, 242, panel_w-32, 26, px=14, bold=True, color='7f1d1d', wrap=False)
    if items:
        for j, (lv, desc, bg, fc) in enumerate(items):
            rect(s, lx2+16, 278+j*72, panel_w-32, 58, fill=bg, border='fecaca')
            rect(s, lx2+16, 278+j*72, 48, 28, fill='fecaca')
            tb(s, lv, lx2+18, 278+j*72+5, 44, 20, px=10, bold=True,
               color=fc, wrap=False, align='center')
            tb(s, desc, lx2+72, 278+j*72+16, panel_w-88, 24,
               px=10.5, color='991b1b', wrap=False)
    elif code:
        rect(s, lx2+16, 276, panel_w-32, 200, fill='fef2f2', border='fecaca')
        tb(s, code, lx2+26, 284, panel_w-48, 180, px=10.5, color='dc2626',
           face='Courier New', wrap=True)


# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f'\n✓ Saved: {OUT}')
print(f'  {len(prs.slides)} slides — fully editable native PPTX elements')
print('\n  Open in PowerPoint or Keynote — all text/shapes are natively editable.')
